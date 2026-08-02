# ==============================
# LIBRERÍAS
# ==============================

# Librerías estándar del sistema
import os  # Para manejar rutas y archivos
import sys  # Para saber cómo se ejecuta el programa (exe o .py)
import copy  # Para clonar objetos (en este caso diapositivas)
import time  # Para pausas del sistema
from datetime import datetime, timedelta  # Manejo de fechas y tiempos
import re

# Librerías externas
import pyperclip  # Para leer texto copiado (portapapeles, ej: WhatsApp)
import win32com.client  # Para controlar PowerPoint desde Python (exportar PDF)
from PIL import Image, ImageEnhance  # Para procesar y mejorar imágenes
from openpyxl import load_workbook, Workbook  # Para leer y crear archivos Excel

# Librerías de PowerPoint
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE


# ==============================
# RUTA BASE DEL PROGRAMA
# ==============================

# Aquí detecto desde dónde se está ejecutando el programa
# Si está compilado (exe), toma la ruta del ejecutable
# Si no, toma la ruta del archivo .py
if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# Rutas principales de archivos
RUTA_EXCEL = os.path.join(BASE_DIR, "incidentes_mes.xlsx")
RUTA_PLANTILLA_8AM = os.path.join(BASE_DIR,"plantilla-8am.pptx")
RUTA_PLANTILLA_2PM = os.path.join(BASE_DIR,"plantilla-2pm.pptx")
RUTA_PLANTILLA_8PM = os.path.join(BASE_DIR,"plantilla-8pm.pptx")

# Rutas donde están las imágenes según turno
CARPETA_IMG_8AM = os.path.join(BASE_DIR,"ImagenesInforme","IMG-8AM")
CARPETA_IMG_2PM = os.path.join(BASE_DIR,"ImagenesInforme","IMG-2PM")
CARPETA_IMG_8PM = os.path.join(BASE_DIR,"ImagenesInforme","IMG-8PM")

# Carpeta donde se guardan los resultados finales
CARPETA_SALIDA = os.path.join(BASE_DIR,"salida")

# Lista de agencias que deberían aparecer sí o sí
AGENCIAS_ESPERADAS = [
    "SUR",
    "CRUE",
    "MOVILIDAD",
    "BOMBEROS",
    "MEBOG",
    "IDIGER"
]

# ==============================
# MESES
# ==============================

# Diccionario para convertir número de mes a texto
MESES = {
    1:"enero",2:"febrero",3:"marzo",4:"abril",
    5:"mayo",6:"junio",7:"julio",8:"agosto",
    9:"septiembre",10:"octubre",11:"noviembre",12:"diciembre"
}

# ==============================
# ESTRUCTURA DEL EXCEL
# ==============================

HOJAS_EXCEL = {
    "Incidentes": [
        "Codigo",
        "Agencia",
        "Tipo",
        "Fecha"
    ],

    "Novedades": [
        "Hora",
        "Agencia",
        "Novedad",
        "Estado"
    ],

    "EstadoGeneral": [
        "TotalAlertas",
        "PorcentajeAutomaticas",
        "InvestigacionesManuales",
        "Critico",
        "Alto",
        "Medio",
        "Bajo",
        "Amenaza",
        "Estacion",
        "Recomendacion"
    ]
}

# ==============================
# VALORES POR DEFECTO
# ==============================

RECOMENDACION_DEFAULT = (
    "Se recomienda reforzar los controles asociados a la plataforma "
    "PremierOne mediante el monitoreo continuo de eventos, la revisión "
    "periódica de los avisos de seguridad (Security Advisories) y la "
    "aplicación de buenas prácticas por parte de los usuarios, "
    "garantizando la resiliencia operativa y el cumplimiento de los "
    "lineamientos de seguridad de la operación del NUSE."
)

ESTADO_GENERAL_DEFAULT = [
    0,                  # TotalAlertas
    0,                  # PorcentajeAutomaticas
    0,                  # InvestigacionesManuales
    0,                  # Critico
    0,                  # Alto
    0,                  # Medio
    0,                  # Bajo
    "",                 # Amenaza
    "",                 # Estacion
    RECOMENDACION_DEFAULT
]

# ==============================
# FORMATO CORPORATIVO TEXTO
# ==============================

# Esta función aplica el formato estándar a los textos del PowerPoint
def aplicar_formato(run, size=16, r=126, g=126, b=126):
    run.font.name = "Calibri"  # Tipo de letra
    run.font.size = Pt(size)  # Tamaño
    run.font.color.rgb = RGBColor(r,g,b)  # Color


# ==============================
# CLONAR DIAPOSITIVA (CON POSICIÓN)
# ==============================

# Esta función copia una diapositiva completa y la ubica en una posición específica
def clonar_diapositiva(prs, slide_index, posicion_destino):

    # Obtengo la diapositiva base
    slide = prs.slides[slide_index]
    layout = slide.slide_layout

    # Creo una nueva diapositiva con el mismo layout
    new_slide = prs.slides.add_slide(layout)

    # Copio todos los elementos (shapes) de la diapositiva original
    for shape in slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # Reorganizo la posición de la nueva diapositiva dentro del archivo
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)

    slide_id = slides[-1]  # tomo la última (la que acabo de crear)

    slide_id_list.remove(slide_id)
    slide_id_list.insert(posicion_destino, slide_id)

    return new_slide

# ==============================
# BUSCAR DIAPOSITIVA POR SHAPE
# ==============================
def obtener_slide_por_shape(prs, nombre_shape):

    for slide in prs.slides:

        for shape in slide.shapes:

            if shape.name == nombre_shape:
                return slide

    return None


# ==============================
# DETECTAR TURNO AUTOMATICO
# ==============================

# Esta función detecta el turno dependiendo de la hora actual
def obtener_turno():

    ahora = datetime.now()
    hora = ahora.hour
    minuto = ahora.minute

    # Convierto la hora a decimal para comparar mejor
    hora_decimal = hora + minuto/60

    # Defino rangos de turno
    if hora_decimal < 8.5:
        return "08_AM", "00:00", "08:00"

    elif hora_decimal < 14.5:
        return "02_PM", "00:00", "14:00"

    elif hora_decimal < 20.5:
        return "08_PM", "00:00", "20:00"

    else:
        return "08_PM", "00:00", "20:00"


# ==============================
# ACTUALIZAR TEXTO DIAPOSITIVA 1
# ==============================

# Esta función actualiza la fecha y el rango horario en la portada
def actualizar_texto(prs, hora_inicio, hora_fin):

    slide = prs.slides[0]  # primera diapositiva

    ahora = datetime.now()
    dia = f"{ahora.day:02d}"
    mes = MESES[ahora.month].capitalize()
    anno = ahora.year

    # Líneas que se van a mostrar
    linea1 = f"Bogotá D.C. {dia} de {mes} del {anno}"
    linea2 = f"Desde las {hora_inicio} hasta las {hora_fin}"

    for shape in slide.shapes:

        # Busco el cuadro de texto correcto
        if shape.name == "txtFecha" and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()  # limpio lo que tenga

            # Primera línea (fecha)
            p1 = tf.paragraphs[0]
            p1.text = linea1
            p1.font.name = "Calibri"
            p1.font.size = Pt(24)
            p1.font.color.rgb = RGBColor(255,255,255)
            p1.alignment = PP_ALIGN.CENTER

            # Segunda línea (horario)
            p2 = tf.add_paragraph()
            p2.text = linea2
            p2.font.name = "Calibri"
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(255,255,255)
            p2.alignment = PP_ALIGN.CENTER

# ============================== 
# ACTUALIZAR TEXTO DIAPOSITIVA 16 (solo 8AM)
# ==============================

# Esta función actualiza los textos de la diapositiva 16,
# pero solo aplica para el informe del turno de las 8AM
def actualizar_estadisticas_8am(prs):

    # Accedo directamente a la diapositiva 16 (índice 15 porque empieza en 0)
    slide = prs.slides[15]

    # ==============================
    # CALCULAR FECHAS
    # ==============================

    # Obtengo la fecha actual
    hoy = datetime.now()

    # Calculo ayer y antier
    ayer = hoy - timedelta(days=1)
    antier = hoy - timedelta(days=2)

    # Formateo las fechas en formato día/mes/año
    fecha_ayer = ayer.strftime("%d/%m/%Y")
    fecha_antier = antier.strftime("%d/%m/%Y")

    # ==============================
    # ARMAR TEXTOS
    # ==============================

    # Texto para estadísticas de un solo día (ayer)
    txt1 = f"Estadística general de llamadas con corte desde las 00:00 AM hasta las 11:59 PM del {fecha_ayer}"

    # Texto para estadísticas acumuladas (antier + ayer)
    txt2 = f"Estadística total llamadas con corte desde las 00:00 AM hasta las 11:59 PM del {fecha_antier} y {fecha_ayer}"

    # Recorro todos los elementos de la diapositiva
    for shape in slide.shapes:

        # Si encuentro el primer cuadro de texto
        if shape.name == "txt-D16-1" and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()  # limpio el contenido anterior

            # Configuración del cuadro
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            # Configuro el párrafo
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            # Agrego el texto
            run = p.add_run()
            run.text = txt1

            # Aplico formato corporativo
            aplicar_formato(run)

        # Segundo cuadro de texto
        elif shape.name == "txt-D16-2" and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()

            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            run = p.add_run()
            run.text = txt2

            aplicar_formato(run)


# ==============================
# LEER INCIDENTES DESDE EXCEL
# ==============================

# Esta función lee los incidentes desde el archivo Excel
def leer_incidentes():

    # Si el archivo no existe, retorno lista vacía
    if not os.path.exists(RUTA_EXCEL):
        return []

    try:
        # Intento abrir el archivo Excel
        wb = load_workbook(RUTA_EXCEL, data_only=True)

    # Si está abierto en Excel, lanza este error
    except PermissionError:
        print(f"❌ ERROR: El archivo Excel está abierto: {RUTA_EXCEL}")
        print("Por favor, ciérralo e intenta de nuevo.")
        time.sleep(5)
        sys.exit(1)

    # Cualquier otro error
    except Exception as e:
        print(f"❌ Error inesperado al abrir el Excel: {e}")
        return []
    
    # A partir de aquí garantizo que el libro se cierre
    try:
        # Si la hoja "Incidentes" no existe, retorno vacío
        if "Incidentes" not in wb.sheetnames:
            return []

        ws = wb["Incidentes"]
        incidentes = []

        # Recorro todas las filas desde la 2 (saltando encabezado)
        for row in ws.iter_rows(min_row=2, values_only=True):

            if not row or len(row) < 4:
                continue

            codigo = row[0]
            agencia = row[1]
            tipo = row[2]
            fecha = row[3]

            if codigo and agencia and tipo:
                incidentes.append({
                    "fecha": fecha,
                    "codigo": str(codigo).strip(),
                    "agencia": str(agencia).strip().upper(),
                    "tipo": str(tipo).strip().lower()
                })

        return incidentes

    finally:
        wb.close()

# ==============================
# CREAR EXCEL DE INCIDENTES SI NO EXISTE
# ==============================

def crear_excel_incidentes():

    archivo_nuevo = False

    # ==============================
    # Abrir o crear Excel
    # ==============================

    if os.path.exists(RUTA_EXCEL):

        try:
            wb = load_workbook(RUTA_EXCEL)

        except PermissionError:
            print("❌ El Excel está abierto.")
            return

    else:
        wb = Workbook()
        archivo_nuevo = True

        # elimina la hoja vacía por defecto
        wb.remove(wb.active)

    # ==============================
    # Validar todas las hojas
    # ==============================

    for nombre_hoja, encabezados in HOJAS_EXCEL.items():

        # --------------------------
        # Crear hoja si no existe
        # --------------------------

        if nombre_hoja not in wb.sheetnames:

            ws = wb.create_sheet(nombre_hoja)

            ws.append(encabezados)

            print(f"✅ Hoja creada: {nombre_hoja}")

            # Datos iniciales únicamente para EstadoGeneral
            if nombre_hoja == "EstadoGeneral":
                ws.append(ESTADO_GENERAL_DEFAULT)

            continue

        # --------------------------
        # Validar encabezados
        # --------------------------

        ws = wb[nombre_hoja]

        actuales = []

        for celda in ws[1]:
            actuales.append("" if celda.value is None else str(celda.value).strip())

        cambio = False

        for columna, esperado in enumerate(encabezados, start=1):

            if columna > len(actuales):

                ws.cell(row=1, column=columna).value = esperado
                cambio = True

            elif actuales[columna-1] != esperado:

                ws.cell(row=1, column=columna).value = esperado
                cambio = True

        # Si EstadoGeneral quedó solo con encabezados
        if nombre_hoja == "EstadoGeneral" and ws.max_row < 2:
            
            ws.append(ESTADO_GENERAL_DEFAULT)

            cambio = True

        if cambio:
            print(f"🔧 Encabezados reparados: {nombre_hoja}")

    wb.save(RUTA_EXCEL)
    wb.close()

    if archivo_nuevo:
        print("✅ Excel creado correctamente.")

# ==============================
# LEER NOVEDADES DESDE EXCEL
# ==============================

# Esta función lee las novedades activas desde el Excel
def leer_novedades():

    # Si no existe, retorno vacío
    if not os.path.exists(RUTA_EXCEL):
        return []

    wb = load_workbook(RUTA_EXCEL, data_only=True)

    # Si no existe la hoja, retorno vacío
    if "Novedades" not in wb.sheetnames:
        return []

    ws = wb["Novedades"]
    novedades = []

    # Recorro las filas
    for row in ws.iter_rows(min_row=2, values_only=True):
        
        # Validación de seguridad (evita errores por filas incompletas)
        if not row or len(row) < 4:
            continue

        fecha = row[0]
        agencia = row[1]
        novedad = row[2]
        estado = row[3]

        # Ignoro registros vacíos
        if novedad is None:
            continue

        if estado is None:
            continue

        # Normalizo el estado
        estado = str(estado).strip().upper()

        # Solo tomo las que estén ACTIVAS
        if estado != "ACTIVA":
            continue

        texto = str(novedad).strip()

        # Si hay texto válido lo agrego
        if texto:
            novedades.append(texto)
    
    wb.close()
    return novedades

# ==============================
# LEER ESTADO GENERAL DESDE EXCEL
# ==============================
def leer_estado_general():

    if not os.path.exists(RUTA_EXCEL):
        return None

    wb = None

    try:
        wb = load_workbook(RUTA_EXCEL, data_only=True)

        if "EstadoGeneral" not in wb.sheetnames:
            print("❌ No existe la hoja EstadoGeneral")
            return None

        ws = wb["EstadoGeneral"]

        # Verifica que exista una fila de datos
        if ws.max_row < 2:
            print("❌ La hoja EstadoGeneral no contiene datos.")
            return None

        # Leer encabezados
        encabezados = [
            str(celda.value).strip() if celda.value else ""
            for celda in ws[1]
        ]

        # Leer datos
        valores = [
            celda.value
            for celda in ws[2]
        ]

        datos = dict(zip(encabezados, valores))

        # Valores por defecto
        datos["TotalAlertas"] = datos.get("TotalAlertas") or 0
        datos["PorcentajeAutomaticas"] = datos.get("PorcentajeAutomaticas") or 0
        datos["InvestigacionesManuales"] = datos.get("InvestigacionesManuales") or 0
        datos["Critico"] = datos.get("Critico") or 0
        datos["Alto"] = datos.get("Alto") or 0
        datos["Medio"] = datos.get("Medio") or 0
        datos["Bajo"] = datos.get("Bajo") or 0
        datos["Amenaza"] = datos.get("Amenaza") or ""
        datos["Estacion"] = datos.get("Estacion") or ""
        recomendacion = datos.get("Recomendacion") or RECOMENDACION_DEFAULT

        recomendacion = (
            str(recomendacion)
            .replace("_x000B_", "\n")
            .replace("\r\n", "\n")
            .replace("\r", "\n")
        )

        datos["Recomendacion"] = recomendacion

        return datos

    except PermissionError:
        print(f"❌ ERROR: El archivo Excel está abierto: {RUTA_EXCEL}")
        return None

    except Exception as e:
        print(f"❌ Error al abrir Excel: {e}")
        return None

    finally:
        if wb is not None:
            wb.close()

# ==============================
# CONTAR TIPOS DE INCIDENTES
# ==============================

# Esta función cuenta cuántos incidentes hay por cada tipo (malware, ransomware u otros)
def contar_tipos_incidente(incidentes):

    # Inicializo contadores
    malware = 0
    ransomware = 0
    otros = 0

    # Recorro todos los incidentes
    for inc in incidentes:

        # Normalizo el texto del tipo para evitar errores de comparación
        tipo = inc["tipo"].lower().strip()

        # Clasifico según el contenido del tipo
        if "malware" in tipo:
            malware += 1
        elif "ransomware" in tipo:
            ransomware += 1
        else:
            otros += 1

    # Retorno los tres conteos
    return malware, ransomware, otros


# ==============================
# DETECTAR AMENAZA MAS REPETIDA
# ==============================

# Esta función identifica cuál código de incidente se repite más
def amenaza_mas_recurrente(incidentes):

    # Diccionario para contar ocurrencias por código
    conteo = {}

    # Recorro los incidentes
    for inc in incidentes:

        codigo = inc["codigo"]

        # Si el código no existe en el diccionario lo inicializo
        if codigo not in conteo:
            conteo[codigo] = 0

        # Sumo una ocurrencia
        conteo[codigo] += 1

    # Si no hay datos, retorno vacío
    if not conteo:
        return None, 0

    # Busco el código con mayor cantidad
    codigo_top = max(conteo, key=conteo.get)
    cantidad = conteo[codigo_top]

    return codigo_top, cantidad


# ==============================
# DETECTAR AGENCIAS AFECTADAS
# ==============================

# Esta función identifica qué agencias están involucradas en los incidentes
def detectar_agencias(incidentes):

    # Mapeo para nombres formateados
    mapa = {
        "SUR": "SUR",
        "CRU": "CRUE",
        "CAD": "CAD",
        "ADM": "ADM",
        
        
        "MOV":"MOVILIDAD",
        "BOM":"BOMBEROS",
        "MEBOG":"MEBOG",
        "IDIGER":"IDIGER",
        
    }

    agencias = set()  # uso set para evitar duplicados

    # Recorro incidentes
    for inc in incidentes:

        codigo = inc["agencia"]

        # Solo agrego si está en el mapa
        if codigo in mapa:
            agencias.add(mapa[codigo])

    # Retorno lista ordenada
    return sorted(list(agencias))


# ==============================
# FORMATEAR LISTA AGENCIAS
# ==============================

# Esta función arma un texto bonito tipo: A, B y C
def formatear_agencias(lista):

    # Si no hay nada
    if not lista:
        return ""

    # Si solo hay una
    if len(lista) == 1:
        return lista[0]

    # Si hay dos
    if len(lista) == 2:
        return f"{lista[0]} y {lista[1]}"

    # Si hay más de dos
    return ", ".join(lista[:-1]) + f" y {lista[-1]}"


# ==============================
# GENERAR TEXTO INCIDENTES D7
# ==============================

# Esta función construye todo el párrafo de la diapositiva 7
def generar_texto_incidentes():
    
    # Leo los incidentes desde el Excel
    incidentes = leer_incidentes()

    # Obtengo estadísticas
    malware, ransomware, otros = contar_tipos_incidente(incidentes)
    codigo_top, cantidad_top = amenaza_mas_recurrente(incidentes)

    tipos = []

    # Armo lista de tipos encontrados
    if malware > 0:
        tipos.append(f"{malware} eventos tipo malware")

    if ransomware > 0:
        tipos.append(f"{ransomware} eventos tipo ransomware")

    if otros > 0:
        tipos.append(f"{otros} eventos de otro tipo")

    # Uno los tipos en un texto
    tipos_txt = " y ".join(tipos)

    frase_top = ""

    # Si hay un evento recurrente significativo lo agrego
    if cantidad_top > 1:

        frase_top = (
            f" Adicionalmente, se identifica que el evento más recurrente fue "
            f"{codigo_top} con {cantidad_top} registros durante el período analizado."
        )

    # Fecha actual
    hoy = datetime.now()
    mes = MESES[hoy.month]
    anno = hoy.year

    # Si no hay incidentes, retorno texto por defecto
    if not incidentes:

        return (
        f"El registro del top de amenazas corresponde al período comprendido "
        f"entre el 1 de {mes} de {anno} y la fecha registrada en el pantallazo. "
        f"En el período indicado no se han identificado incidentes de seguridad "
        f"asociados a software malicioso en las plataformas monitoreadas. "
        f"Es importante reforzar los controles de seguridad asociados al software malicioso, "
        f"tales como: monitoreo periódico de eventos de seguridad, actualización continua "
        f"de herramientas antivirus, respaldo de las soluciones de seguridad y programas "
        f"de concientización a los usuarios sobre phishing, descargas no autorizadas y "
        f"buenas prácticas de seguridad de la información. La consolidación de estas "
        f"medidas permitirá reducir el riesgo de incidentes, mantener la resiliencia "
        f"operativa y garantizar el cumplimiento de los lineamientos de la operación."
        f"{frase_top}"
        )

    # Total de incidentes
    total = len(incidentes)

    # Detecto agencias involucradas
    agencias = detectar_agencias(incidentes)
    agencias_txt = formatear_agencias(agencias)

    codigos = []

    # Armo lista de códigos con tipo
    for inc in incidentes:
        codigo = inc["codigo"]
        tipo = inc["tipo"].capitalize()
        fecha = inc["fecha"]

        # Si Excel devuelve un objeto datetime
        if isinstance(fecha, datetime):
            fecha = fecha.strftime("%d/%m/%Y")
        else:
            fecha = str(fecha)

        codigos.append(f"({codigo} {tipo} ({fecha}))")

    # Formateo la lista de códigos
    if len(codigos) == 1:
        lista = codigos[0]
    else:
        lista = ", ".join(codigos[:-1]) + " y " + codigos[-1]

    # Texto final completo
    texto = (
        f"El registro del top de amenazas corresponde al período comprendido "
        f"entre el 1 de {mes} de {anno} y la fecha registrada en el pantallazo. "
        f"En el período indicado se han identificado {total} incidentes de seguridad "
        f"en las agencias {agencias_txt}, correspondientes a {tipos_txt}, "
        f"específicamente: {lista}. "
        f"Es importante reforzar los controles de seguridad asociados al software malicioso, "
        f"tales como: monitoreo periódico de eventos de seguridad, actualización continua "
        f"de herramientas antivirus, respaldo de las soluciones de seguridad y programas "
        f"de concientización a los usuarios sobre phishing, descargas no autorizadas y "
        f"buenas prácticas de seguridad de la información. La consolidación de estas "
        f"medidas permitirá reducir el riesgo de incidentes, mantener la resiliencia "
        f"operativa y garantizar el cumplimiento de los lineamientos de la operación."
        f"{frase_top}"
    )

    return texto

#====================================
# Generar el texto de estado General.
#====================================

def generar_texto_estado_general(datos):

    hoy = datetime.now()
    mes = MESES[hoy.month]
    anno = hoy.year

    total_criticidad = (
        datos.get("Bajo", 0)
        + datos.get("Medio", 0)
        + datos.get("Alto", 0)
        + datos.get("Critico", 0)
    )

    if total_criticidad == 0:

        segundo_parrafo = (
            "No se identificaron amenazas creíbles ni casos de seguridad "
            "(P1-P4) en el marco de tiempo evaluado y no existen alertas "
            "activas en ningún nivel de criticidad."
        )

    else:

        segundo_parrafo = (
            f"Durante el período evaluado se identificaron amenazas de "
            f"{datos.get('Amenaza','Malware')} creíbles en la estación "
            f"{datos.get('Estacion','N/D')} (P1-P4). Asimismo, se mantienen "
            "alertas activas con la siguiente clasificación de criticidad."
        )

    return [

        f"Durante {mes} de {anno}, con el monitoreo SOC activo, "
        f"se generaron {datos.get('TotalAlertas',0)} alertas de seguridad, "
        f"de las cuales el {datos.get('PorcentajeAutomaticas',0)}% fueron "
        f"resueltas de forma automatizada y "
        f"{datos.get('InvestigacionesManuales',0)} requirieron investigación manual del equipo SOC.",

        segundo_parrafo,

        f"Bajo: ({datos.get('Bajo',0)})",
        f"Medio: ({datos.get('Medio',0)})",
        f"Alto: ({datos.get('Alto',0)})",
        f"Crítico: ({datos.get('Critico',0)})",

        datos.get("Recomendacion", ""),

        "P1 – Low (Bajo)",
        "P2 – Medium (Medio)",
        "P3 – High (Alto)",
        "P4 – Critical (Crítico)"
    ]
# ==============================
# ACTUALIZAR DIAPOSITIVA 7-1
# ==============================
def actualizar_estado_general(prs):

    slide = obtener_slide_por_shape(prs, "txt-D7-1")

    if slide is None:
        print("❌ No se encontró la diapositiva Estado General.")
        return

    # Leer la información desde Excel
    datos = leer_estado_general()

    if datos is None:
        print("❌ No fue posible leer la hoja EstadoGeneral.")
        return

    # Generar el texto utilizando los datos leídos
    texto = generar_texto_estado_general(datos)

    # Determinar tamaño de letra según la cantidad de alertas
    total_criticidad = (
        datos.get("Bajo", 0)
        + datos.get("Medio", 0)
        + datos.get("Alto", 0)
        + datos.get("Critico", 0)
    )

    tamano = 13 if total_criticidad == 0 else 12

    # Buscar el cuadro de texto
    for shape in slide.shapes:

        if shape.name == "txt-D7-1" and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()

            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = MSO_ANCHOR.TOP

            for i, linea in enumerate(texto):

                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.space_before = Pt(0)

                # Un poco de separación entre párrafos
                if (
                    linea.startswith("P")
                    or linea.startswith("Bajo")
                    or linea.startswith("Medio")
                    or linea.startswith("Alto")
                    or linea.startswith("Crítico")
                ):
                    p.space_after = Pt(2)
                    p.alignment = PP_ALIGN.LEFT
                else:
                    p.space_after = Pt(6)
                    p.alignment = PP_ALIGN.JUSTIFY

                run = p.add_run()
                run.text = linea.replace("\n", "\v")

                aplicar_formato(run, size=tamano)

            # Ya encontró el cuadro de texto
            break

# ==============================
# ACTUALIZAR DIAPOSITIVA 7
# ==============================
def actualizar_diapositiva_7(prs):

    slide = obtener_slide_por_shape(prs, "txt-D7")

    if slide is None:
        print("❌ No se encontró la diapositiva de Incidentes.")
        return
    
    texto = generar_texto_incidentes()  # Genera automáticamente el texto de incidentes

    for shape in slide.shapes:  # Recorre todos los elementos de la diapositiva

        # Busca el cuadro de texto específico llamado "txt-D7"
        if shape.name == "txt-D7" and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()  # Borra contenido anterior

            # CONFIGURACIÓN DEL CUADRO DE TEXTO
            tf.word_wrap = True  # Permite salto automático de línea
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Ajusta el texto al tamaño del cuadro

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.JUSTIFY   # Texto justificado (alineado a ambos lados)

            run = p.add_run()
            run.text = texto  # Inserta el texto generado

            aplicar_formato(run, size=14)   # Aplica formato corporativo con tamaño 14


# ==============================
# ACTUALIZAR TITULO NOVEDADES
# ==============================
def actualizar_titulo_novedades(slide, numero):

    for shape in slide.shapes:  # Recorre todos los elementos de la diapositiva

        # Busca el título que empieza con "titulo-D"
        if shape.name.startswith("titulo-D") and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()  # Limpia el contenido

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT  # Alineación a la izquierda

            run = p.add_run()
            run.text = f"{numero}. Novedades"  # Inserta el número dinámico del título

            # Formato del título
            run.font.name = "Calibri"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(27, 95, 167)


# ==============================
# CALCULAR TAMAÑO DINÁMICO TEXTO
# ==============================
"""def calcular_tamano_bloque(novedades):

    # Obtiene la novedad más larga de la diapositiva
    mayor = max(len(n) for n in novedades)

    # Tamaños según la longitud del texto más largo
    if mayor <= 120:
        return 18

    elif mayor <= 220:
        return 16

    elif mayor <= 350:
        return 14

    else:
        return 12"""

def calcular_tamano_bloque(novedades):

    total = sum(len(n) for n in novedades)

    if total < 300:
        return 18
    elif total < 600:
        return 16
    elif total < 900:
        return 14
    else:
        return 12
    

# ==============================
# ACTUALIZAR DIAPOSITIVA 26
# ==============================
def actualizar_diapositiva_26(prs):

    slide_index = 26  # Índice base (diapositiva 26)

    # Validación: verifica que la plantilla tenga suficientes diapositivas
    if len(prs.slides) <= slide_index:
        print("❌ La plantilla no tiene suficientes diapositivas")
        return

    novedades = leer_novedades()  # Obtiene novedades desde Excel

    # Si no hay novedades, coloca mensaje por defecto
    if not novedades:
        novedades = [
            "Durante el periodo del informe no se registran novedades relevantes en la operación."
        ]

    MAX_POR_SLIDE = 3  # Máximo de novedades por diapositiva
    BASE_TITULO = 16   # Número base para numerar títulos

    # Divide las novedades en bloques de máximo 3
    bloques = [
        novedades[i:i + MAX_POR_SLIDE]
        for i in range(0, len(novedades), MAX_POR_SLIDE)
    ]

    tamano = calcular_tamano_bloque(novedades)

    for i, bloque in enumerate(bloques):

        # ==============================
        # CREAR / OBTENER SLIDE
        # ==============================
        if i == 0:
            slide = prs.slides[slide_index]  # Usa la diapositiva base
        else:
            posicion = len(prs.slides) - 2  # Inserta antes de la última
            slide = clonar_diapositiva(prs, slide_index, posicion)  # Clona la base

        # ==============================
        # RENOMBRAR SHAPES (CLAVE)
        # ==============================
        for shp in slide.shapes:

            if shp.name == "txt-D26":
                shp.name = f"txt-D{26 + i}"  # Renombra dinámicamente

            if shp.name == "titulo-D26":
                shp.name = f"titulo-D{26 + i}"

        # ==============================
        # ACTUALIZAR TÍTULO
        # ==============================
        actualizar_titulo_novedades(slide, BASE_TITULO + i)

        # ==============================
        # ESCRIBIR NOVEDADES
        # ==============================
        nombre_shape = f"txt-D{26 + i}"  # Nombre dinámico del cuadro de texto

        for shape in slide.shapes:

            if shape.name == nombre_shape and shape.has_text_frame:

                tf = shape.text_frame
                tf.clear()

                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.NONE  # No cambia tamaño del cuadro
                tf.vertical_anchor = MSO_ANCHOR.TOP

                # Inserta cada novedad como viñeta
                
                

                primera_linea = True

                for novedad in bloque:

                    lineas = novedad.split("\n")

                    for num_linea, linea in enumerate(lineas):

                        if primera_linea:
                            p = tf.paragraphs[0]
                            primera_linea = False
                        else:
                            p = tf.add_paragraph()

                        p.alignment = PP_ALIGN.LEFT

                        # Solo deja un pequeño espacio al terminar una novedad
                        if num_linea == len(lineas) - 1:
                            p.space_after = Pt(10)
                        else:
                            p.space_after = Pt(0)

                        run = p.add_run()

                        if num_linea == 0:
                            run.text = f"• {linea}"
                        else:
                            run.text = linea

                        run.font.name = "Calibri"
                        run.font.size = Pt(tamano)
                        run.font.color.rgb = RGBColor(156, 156, 156)


# ==============================
# FORMATEAR NOMBRES
# ==============================
def formatear_nombre(nombre):

    if not nombre:
        return ""

    nombre = str(nombre)

    # Normaliza espacios especiales (WhatsApp)
    nombre = nombre.replace("\xa0", " ")

    # Elimina cualquier cantidad de :, ., -, > al inicio
    nombre = re.sub(r'^[\s:.\-<>]+', '', nombre)

    # Elimina espacios repetidos
    nombre = " ".join(nombre.split())

    partes = [p.capitalize() for p in nombre.lower().split()]

    return " ".join(partes)
# ===================

def normalizar_agencia(texto):

    if not texto:
        return ""

    texto = str(texto).upper()

    texto = texto.replace("\xa0", " ")

    texto = re.sub(r'^[^A-ZÁÉÍÓÚÑ]+', '', texto)

    texto = re.sub(r'[^A-ZÁÉÍÓÚÑ ]', '', texto)

    texto = " ".join(texto.split())

    return texto

#==================================
# ======================================
# EXTRAER AGENCIA
# ======================================
def extraer_agencia(linea):

    m = re.search(
        r"UBICACI[ÓO]N\s*:*\s*:*\s*(.+?)\s*$",
        linea,
        re.IGNORECASE
    )

    if not m:
        return None

    return normalizar_agencia(m.group(1))


# ======================================
# EXTRAER JEFE
# ======================================
def extraer_jefe(linea):

    linea = linea.strip()

    m = re.search(r"JEFE\s+DE\s+SALA\s*:*\s*:*\s*(.+?)\s*$", linea, re.I)

    if m:
        return formatear_nombre(m.group(1))

    return None


# ======================================
# EXTRAER FORMATO CORTO
# ======================================
def extraer_formato_corto(linea):

    linea = linea.strip()

    patrones = [
        r"^(.+?)\s*-\s*(.+)$",
        r"^(.+?)\s*-->\s*:*\s*(.+)$",
        r"^(.+?)\s*:+\s*:*\s*(.+)$",
    ]

    for patron in patrones:

        m = re.match(patron, linea)

        if m:

            agencia = normalizar_agencia(m.group(1))
            nombre = formatear_nombre(m.group(2))

            return agencia, nombre

    return None, None
#==================================


# ==============================
# EXTRAER NOMBRES WHATSAPP
# ==============================
def obtener_disponibles_portapapeles():

    disponibles = {}

    texto = pyperclip.paste()
    lineas = texto.splitlines()

    agencia = None

    for linea in lineas:

        linea = linea.strip()

        if not linea:
            continue

        # -------------------------
        # UBICACIÓN
        # -------------------------
        ag = extraer_agencia(linea)

        if ag:
            agencia = ag
            continue

        # -------------------------
        # JEFE DE SALA
        # -------------------------
        jefe = extraer_jefe(linea)

        if jefe:

            if agencia:
                disponibles[agencia] = jefe

            agencia = None
            continue

        # -------------------------
        # FORMATO CORTO
        # -------------------------
        ag, jefe = extraer_formato_corto(linea)

        if ag and jefe:
            disponibles[ag] = jefe

    print("\n===== DISPONIBLES =====")
    for a, n in disponibles.items():
        print(a, "->", n)

    return disponibles

# ==========================================================
# AJUSTAR IMAGEN CON MARGEN Y MEJORAR CALIDAD (PRO)
# ==========================================================
def ajustar_imagen_con_margen(ruta, marco, color_fondo=(255, 255, 255)):

    try:
        # Abre la imagen
        with Image.open(ruta) as img_raw:
            img = img_raw.convert("RGB")  # Convierte a formato RGB
            img_w, img_h = img.size  # Dimensiones originales

            # Dimensiones del marco en PowerPoint
            marco_w = int(marco.width)
            marco_h = int(marco.height)

            # Limita tamaño máximo para evitar imágenes gigantes
            MAX_PIX = 4000
            escala_w = min(marco_w, MAX_PIX)
            escala_h = min(marco_h, MAX_PIX)

            # Calcula proporciones
            ratio_img = img_w / img_h
            ratio_marco = escala_w / escala_h

            # Ajusta tamaño manteniendo proporción
            if ratio_img > ratio_marco:
                new_w = escala_w
                new_h = int(escala_w / ratio_img)
            else:
                new_h = escala_h
                new_w = int(escala_h * ratio_img)

            # Redimensiona con alta calidad
            img_resized = img.resize((new_w, new_h), Image.LANCZOS)

            # Mejora nitidez
            sharp = ImageEnhance.Sharpness(img_resized)
            img_resized = sharp.enhance(1.2)

            # Crea fondo blanco
            fondo = Image.new("RGB", (new_w, new_h), color_fondo)

            # Centra la imagen en el fondo
            pos_x = (new_w - img_resized.width) // 2
            pos_y = (new_h - img_resized.height) // 2
            fondo.paste(img_resized, (pos_x, pos_y))

            # Genera nombre de archivo temporal
            nombre_base = os.path.splitext(os.path.basename(ruta))[0]
            ruta_temp = os.path.join(os.path.dirname(ruta), f"{nombre_base}_PRO.png")

            # 🔥 Guarda siempre la imagen procesada
            fondo.save(ruta_temp, format="PNG", dpi=(300, 300))

            # Retorna ruta y posición para insertar en PowerPoint
            return ruta_temp, marco.left, marco.top, marco.width, marco.height

    except Exception as e:
        print(f"❌ Error procesando imagen {ruta}: {e}")

        # 🔥 Si falla, usa la imagen original para no romper el flujo
        return ruta, marco.left, marco.top, marco.width, marco.height


# ==============================
# PROCESAR IMAGENES (MEJORADO)
# ==============================
def procesar_imagenes(prs, CARPETA_IMAGENES):

    # 🔥 Verifica que la carpeta exista
    if not os.path.exists(CARPETA_IMAGENES):
        print("⚠ Carpeta de imágenes no encontrada:", CARPETA_IMAGENES)
        return 0

    total = 0  # Contador de imágenes insertadas
    extensiones = (".png",".jpg",".jpeg")

    for archivo in os.listdir(CARPETA_IMAGENES):
        
        nombre = archivo.lower()

        # Ignora las imágenes temporales generadas por el programa
        if archivo.endswith("_PRO.png"):
            continue
        

        # Filtra solo imágenes válidas
        if not archivo.lower().endswith(extensiones):
            continue

        # Solo procesa archivos que tengan "-d" (ej: imagen-d7.png)
        if "-d" not in nombre:
            continue

        ruta = os.path.join(CARPETA_IMAGENES, archivo)

        # Extrae el código (ej: D7)
        codigo = nombre.split("-d")[1].split(".")[0].upper()
        codigo = f"D{codigo}"

        nombre_marco = f"imgMarco-{codigo}"  # Nombre del marco en PPT

        marco_encontrado = False

        # Busca el slide que tenga ese marco
        for slide in prs.slides:

            if marco_encontrado:
                break

            for shape in slide.shapes:

                # Si encuentra el marco correspondiente
                if shape.name.upper().startswith(nombre_marco.upper()):

                    marco_encontrado = True

                    # Elimina imagen anterior si ya existe
                    for s in slide.shapes:
                        if s.name == f"imgAuto-{codigo}":
                            slide.shapes._spTree.remove(s._element)

                    # Ajusta la imagen (tamaño, calidad, fondo)
                    ruta_temp, left, top, new_w, new_h = ajustar_imagen_con_margen(
                        ruta, shape, color_fondo=(255,255,255)
                    )

                    # Inserta la imagen en PowerPoint
                    pic = slide.shapes.add_picture(
                        ruta_temp,
                        left,
                        top,
                        width=new_w,
                        height=new_h
                    )
                   
                    pic.name = f"imgAuto-{codigo}"  # Nombre interno de la imagen

                    total += 1

                    print("Imagen insertada:", archivo)

                    break

        # Si no encontró marco para esa imagen
        if not marco_encontrado:
            print("⚠ No se encontró marco para:", archivo)

    return total  # Retorna cuántas imágenes insertó

# ==============================
# LIMPIAR IMAGENES TEMPORALES
# ==============================
def limpiar_imagenes_temporales(carpeta):

    if not os.path.exists(carpeta):
        return

    for archivo in os.listdir(carpeta):

        if archivo.endswith("_PRO.png"):

            ruta = os.path.join(carpeta, archivo)

            try:
                os.remove(ruta)
                print("🧹 Eliminado:", archivo)
            except Exception as e:
                print(f"⚠ No se pudo eliminar {archivo}: {e}")


# ==============================
# LIMPIAR RANGO
# ==============================
def limpiar_rango(nombre):

    if not nombre:
        return ""

    nombre = str(nombre).lower()

    # Elimina signos de puntuación
    nombre = re.sub(r"[.,:;]", " ", nombre)


    # Elimina espacios repetidos
    nombre = " ".join(nombre.split())

    rangos = {
        "mayor", "my","m.y",
        "teniente", "tn", "tte","t.te",
        "capitan", "cap", "cpt",
        "coronel", "cr", "cor",
        "subteniente", "st", "si",
        "sargento", "sg", "sgto", "sgt",
        "intendente", "int", "it",
        "patrullero", "pt",
        "oficial", "of",
        "doctor", "doct", "dr"
    }

    partes = [p for p in nombre.split() if p not in rangos]

    return " ".join(partes).title()

# ==============================
# INSERTAR NOMBRES
# ==============================
def insertar_nombres(prs, disponibles):

    # Selecciona la diapositiva 3 (índice 2)
    slide = prs.slides[2]

    # Recorre todos los elementos (shapes) de la diapositiva
    for shape in slide.shapes:

        # Obtiene el nombre del shape (se usa como identificador de agencia)
        agencia = normalizar_agencia(shape.name)

        # Verifica que el shape corresponda a una agencia esperada y tenga texto
        if agencia in AGENCIAS_ESPERADAS and shape.has_text_frame:

            # Si hay nombre disponible para esa agencia
            if agencia in disponibles:
                # Limpia el rango (ej: Mayor, Teniente, etc.)
                nombre = limpiar_rango(disponibles[agencia])
            else:
                # Si no hay datos, coloca texto por defecto
                nombre = "Sin reporte"

            # Accede al cuadro de texto del shape
            tf = shape.text_frame
            tf.clear()  # Limpia contenido anterior

            # Toma el primer párrafo
            p = tf.paragraphs[0]

            # Agrega un nuevo texto (run)
            run = p.add_run()
            run.text = f"{agencia}: {nombre}"

            # Aplica formato corporativo (fuente, tamaño, color)
            aplicar_formato(run)

            # Evita salto automático de línea
            tf.word_wrap = False

            # Ajusta el ancho del cuadro de texto
            shape.width = Inches(4.5)


# ==============================
# EXPORTAR A PDF
# ==============================
def exportar_pdf(ruta_pptx, ruta_pdf):
    powerpoint = None
    try:
        # Abre PowerPoint mediante COM (solo funciona en Windows)
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")

        # Abre la presentación sin mostrar ventana
        presentation = powerpoint.Presentations.Open(os.path.abspath(ruta_pptx), WithWindow=False)

        # Guarda como PDF (formato 32)
        presentation.SaveAs(os.path.abspath(ruta_pdf), 32)

        # Cierra la presentación
        presentation.Close()

    except Exception as e:
        # Manejo de errores en caso de fallo
        print("⚠ No se pudo exportar a PDF.")
        print("Posibles causas:")
        print("- PowerPoint no está instalado")
        print("- Error en COM de Windows")
        print("- Archivo PPT abierto")
        print("Detalle:", e)

    finally:
        # Cierra la aplicación de PowerPoint si fue abierta
        if powerpoint:
            powerpoint.Quit()


# ==============================
# GENERAR INFORME
# ==============================
def generar_informe(turno_manual=None):

    # Asegura que exista el Excel de incidentes
    crear_excel_incidentes()

    # ==============================
    # 1. Definir el turno
    # ==============================
    if turno_manual == "mañana":
        turno, hora_inicio, hora_fin = "08_AM", "00:00", "08:00"
    elif turno_manual == "tarde":
        turno, hora_inicio, hora_fin = "02_PM", "00:00", "14:00"
    elif turno_manual == "noche":
        turno, hora_inicio, hora_fin = "08_PM", "00:00", "20:00"
    else:
        # Si no se especifica, detecta automáticamente
        turno, hora_inicio, hora_fin = obtener_turno()

    # ==============================
    # 2. Seleccionar rutas según el turno
    # ==============================
    if turno == "08_AM":
        ruta_plantilla = RUTA_PLANTILLA_8AM
        CARPETA_IMAGENES = CARPETA_IMG_8AM
    elif turno == "02_PM":
        ruta_plantilla = RUTA_PLANTILLA_2PM
        CARPETA_IMAGENES = CARPETA_IMG_2PM
    else:
        ruta_plantilla = RUTA_PLANTILLA_8PM
        CARPETA_IMAGENES = CARPETA_IMG_8PM

    # ==============================
    # Validación: verificar que exista la plantilla
    # ==============================
    if not os.path.exists(ruta_plantilla):
        print(f"❌ ERROR: No se encontró la plantilla en: {ruta_plantilla}")
        return  # Detiene la ejecución si no existe

    # ==============================
    # 3. Abrir la presentación
    # ==============================
    try:
        prs = Presentation(ruta_plantilla)
    except Exception as e:
        print(f"❌ Error al abrir la presentación: {e}")
        return

    # ==============================
    # Actualizaciones de contenido
    # ==============================
    actualizar_texto(prs, hora_inicio, hora_fin)   # Diapositiva 1
    actualizar_diapositiva_7(prs)                 # Incidentes
    actualizar_estado_general(prs)               # Estado General 
    actualizar_diapositiva_26(prs)                # Novedades

    # Obtiene nombres desde el portapapeles (WhatsApp)
    disponibles = obtener_disponibles_portapapeles()

    # Valida si faltan agencias
    for agencia in AGENCIAS_ESPERADAS:
        if agencia not in disponibles:
            print("⚠ Falta jefe de sala:", agencia)

    # Inserta los nombres en la diapositiva correspondiente
    insertar_nombres(prs, disponibles)

    # ==============================
    # PROCESAR IMÁGENES
    # ==============================
    total = procesar_imagenes(prs, CARPETA_IMAGENES)

    # Solo para turno de la mañana (8AM)
    if turno == "08_AM":
        actualizar_estadisticas_8am(prs)

    # ==============================
    # GENERAR NOMBRE DE ARCHIVO
    # ==============================
    ahora = datetime.now()
    dia = ahora.day
    mes = MESES[ahora.month]

    nombre_archivo = f"Reporte Seguimiento Operación NUSE 123 {dia} de {mes} {turno}"

    # ==============================
    # CREAR CARPETAS DE SALIDA
    # ==============================
    carpeta_informe = os.path.join(CARPETA_SALIDA, "informe")
    carpeta_pdf = os.path.join(CARPETA_SALIDA, f"{mes}-pdf")

    os.makedirs(carpeta_informe, exist_ok=True)
    os.makedirs(carpeta_pdf, exist_ok=True)

    # Rutas finales
    ruta_pptx = os.path.join(carpeta_informe, f"{nombre_archivo}.pptx")
    ruta_pdf = os.path.join(carpeta_pdf, f"{nombre_archivo}.pdf")

    # ==============================
    # GUARDAR PRESENTACIÓN
    # ==============================
    try:
        prs.save(ruta_pptx)
        print("Guardando PPT:", ruta_pptx)
    except PermissionError:
        print(f"❌ ERROR: No se pudo guardar el PPT. ¿Está abierto el archivo {nombre_archivo}.pptx?")
        return

#    time.sleep(3)

   # ==============================
   # EXPORTAR A PDF
   # ==============================
    exportar_pdf(ruta_pptx, ruta_pdf)

    # Espera hasta 5 segundos a que PowerPoint termine de crear el PDF
    for _ in range(10):
        if os.path.exists(ruta_pdf):
            break
        time.sleep(0.5)  # Da tiempo a que PowerPoint termine la exportación
    
    # ==============================
    # VALIDAR CREACIÓN DEL PDF
    # ==============================
    if not os.path.exists(ruta_pdf):
        print("\n⚠ No se generó el PDF")
        print("Revisa PowerPoint o posibles errores de exportación")
        print("El archivo PPT sí fue generado correctamente:")
        print("PPT:", ruta_pptx)

        # Abre el PPT si falla el PDF
        os.startfile(ruta_pptx)

        # Elimina las imágenes temporales
        limpiar_imagenes_temporales(CARPETA_IMAGENES)

        return

    # ==============================
    # MENSAJE FINAL
    # ==============================
    print("\nInforme generado correctamente")
    print("Imágenes insertadas:", total)
    print("PPT:", ruta_pptx)
    print("PDF:", ruta_pdf)

    # Abre automáticamente el PDF
    os.startfile(ruta_pdf)
    # Elimina las imágenes temporales
    limpiar_imagenes_temporales(CARPETA_IMAGENES)

# ==============================
# EJECUTAR SOLO EN CONSOLA
# ==============================
if __name__ == "__main__":

    # Indica que se ejecuta desde consola
    print("Modo consola")

    # Detecta turno actual
    turno, h1, h2 = obtener_turno()

    # Ejecuta generación del informe
    generar_informe()

"""
if __name__ == "__main__":
    print("Modo consola")

    generar_informe("mañana")
"""
    # Autor - (Luis Rojas Rincon)