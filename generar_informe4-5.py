import os
import sys
import time
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from PIL import Image
import pyperclip
import win32com.client

# ==============================
# CONFIGURACION
# ==============================

#RUTA_PLANTILLA = r"C:/Users/LuisAlvaroRojasRinco/Documents/Luis-R/Autom/AutInforme-py/plantilla-2pm.pptx"
#CARPETA_IMAGENES = r"C:/Users/LuisAlvaroRojasRinco/Documents/Luis-R/Autom/AutInforme-py/ImagenesInforme/IMG-2PM"
#CARPETA_SALIDA = r"C:/Users/LuisAlvaroRojasRinco/Documents/Luis-R/Autom/AutInforme-py/salida"

# ==============================
# RUTA BASE DEL PROGRAMA
# ==============================
if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))
RUTA_PLANTILLA_8AM = os.path.join(BASE_DIR,"plantilla-8am.pptx")
RUTA_PLANTILLA_2PM = os.path.join(BASE_DIR,"plantilla-2pm.pptx")
RUTA_PLANTILLA_8PM = os.path.join(BASE_DIR,"plantilla-8pm.pptx")

#Rutas de imágenes automáticas
CARPETA_IMG_8AM = os.path.join(BASE_DIR,"ImagenesInforme","IMG-8AM")
CARPETA_IMG_2PM = os.path.join(BASE_DIR,"ImagenesInforme","IMG-2PM")
CARPETA_IMG_8PM = os.path.join(BASE_DIR,"ImagenesInforme","IMG-8PM")

#Carpeta de salida
CARPETA_SALIDA = os.path.join(BASE_DIR,"salida")

AGENCIAS_ESPERADAS = [
    "SUR",
    "CRUE",
    "MOVILIDAD",
    "BOMBEROS",
    "MEBOG",
    "IDIGER"
]

# ==============================
# MESES
# ==============================
MESES = {
    1:"enero",2:"febrero",3:"marzo",4:"abril",
    5:"mayo",6:"junio",7:"julio",8:"agosto",
    9:"septiembre",10:"octubre",11:"noviembre",12:"diciembre"
}

# ==============================
# DETECTAR TURNO
# ==============================
def obtener_turno():

    ahora = datetime.now()
    hora = ahora.hour
    minuto = ahora.minute

    hora_decimal = hora + minuto/60

    # Informe de 8 AM
    if hora_decimal < 8.5:
        return "08_AM", "00:00", "08:00"

    # Informe de 2 PM
    elif hora_decimal < 14.5:
        return "02_PM", "00:00", "14:00"

    # Informe de 8 PM
    elif hora_decimal < 20.5:
        return "08_PM", "00:00", "20:00"

    # Después de 20:30 sigue siendo informe de 8 PM
    else:
        return "08_PM", "00:00", "20:00"
# ==============================
# ACTUALIZAR TEXTO DIAPOSITIVA 1
# ==============================
def actualizar_texto(prs, hora_inicio, hora_fin):

    slide = prs.slides[0]

    ahora = datetime.now()
    dia = f"{ahora.day:02d}"
    mes = MESES[ahora.month].capitalize()
    anio = ahora.year

    linea1 = f"Bogotá D.C. {dia} de {mes} del {anio}"
    linea2 = f"Desde las {hora_inicio} hasta las {hora_fin}"

    for shape in slide.shapes:

        if shape.name == "txtFecha" and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()

            p1 = tf.paragraphs[0]
            p1.text = linea1
            p1.font.name = "Calibri"
            p1.font.size = Pt(24)
            p1.font.color.rgb = RGBColor(255,255,255)
            p1.alignment = PP_ALIGN.CENTER

            p2 = tf.add_paragraph()
            p2.text = linea2
            p2.font.name = "Calibri"
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(255,255,255)
            p2.alignment = PP_ALIGN.CENTER

# ==============================
# FORMATEAR NOMBRES
# ==============================
def formatear_nombre(nombre):

    partes = nombre.strip().lower().split()
    partes = [p.capitalize() for p in partes]

    return " ".join(partes)

# ==============================
# EXTRAER NOMBRES WHATSAPP
# ==============================
def obtener_disponibles_portapapeles():

    disponibles = {}

    texto = pyperclip.paste()
    lineas = texto.splitlines()

    agencia = ""
    nombre = ""

    for linea in lineas:

        linea = linea.strip()

        if linea.upper().startswith("UBICACIÓN:"):

            agencia = linea.split(":")[1].strip().upper()

        elif linea.upper().startswith("JEFE DE SALA:"):

            nombre = formatear_nombre(linea.split(":")[1].strip())

            if agencia and nombre:

                disponibles[agencia] = nombre
                agencia = ""
                nombre = ""

        elif "-" in linea:

            partes = linea.split("-")

            if len(partes) == 2:

                agencia = partes[0].strip().upper()
                nombre = formatear_nombre(partes[1].strip())

                disponibles[agencia] = nombre

    print("Disponibles detectados:", disponibles)

    return disponibles

# ==============================
# AJUSTAR IMAGEN
# ==============================
def ajustar_imagen(ruta, marco):

    img = Image.open(ruta)
    img_w, img_h = img.size

    marco_w = marco.width
    marco_h = marco.height

    ratio_img = img_w / img_h
    ratio_marco = marco_w / marco_h

    if ratio_img > ratio_marco:

        new_w = marco_w
        new_h = marco_w / ratio_img

    else:

        new_h = marco_h
        new_w = marco_h * ratio_img

    left = marco.left + (marco_w - new_w) / 2
    top = marco.top + (marco_h - new_h) / 2

    return left, top, new_w, new_h

# ==============================
# PROCESAR IMAGENES (MEJORADO)
# ==============================
def procesar_imagenes(prs):

    total = 0
    extensiones = (".png",".jpg",".jpeg")

    for archivo in os.listdir(CARPETA_IMAGENES):

        if not archivo.lower().endswith(extensiones):
            continue

        nombre = archivo.lower()

        if "-d" not in nombre:
            continue

        ruta = os.path.join(CARPETA_IMAGENES, archivo)

        codigo = nombre.split("-d")[1].split(".")[0].upper()
        codigo = f"D{codigo}"

        nombre_marco = f"imgMarco-{codigo}"

        marco_encontrado = False

        for slide in prs.slides:

            if marco_encontrado:
                break

            for shape in slide.shapes:

                if shape.name.upper().startswith(nombre_marco.upper()):

                    marco_encontrado = True

                    for s in slide.shapes:

                        if s.name == f"imgAuto-{codigo}":

                            slide.shapes._spTree.remove(s._element)

                    left, top, new_w, new_h = ajustar_imagen(ruta, shape)

                    pic = slide.shapes.add_picture(
                        ruta,
                        left,
                        top,
                        width=new_w,
                        height=new_h
                    )

                    pic.name = f"imgAuto-{codigo}"

                    total += 1

                    print("Imagen insertada:", archivo)

                    break

        if not marco_encontrado:

            print("⚠ No se encontró marco para:", archivo)

    return total

# ==============================
# LIMPIAR RANGO
# ==============================
def limpiar_rango(nombre):

    rangos = ["mayor","teniente","sargento","capitan","coronel","subteniente"]

    partes = nombre.lower().split()

    partes = [p for p in partes if p not in rangos]

    return " ".join(partes).title()

# ==============================
# INSERTAR NOMBRES
# ==============================
def insertar_nombres(prs, disponibles):

    slide = prs.slides[2]

    for shape in slide.shapes:

        agencia = shape.name.strip().upper()

        if agencia in AGENCIAS_ESPERADAS and shape.has_text_frame:

            if agencia in disponibles:
                nombre = limpiar_rango(disponibles[agencia])
            else:
                nombre = "Sin reporte"

            tf = shape.text_frame
            tf.clear()

            p = tf.paragraphs[0]

            p.text = f"{agencia}: {nombre}"

            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(125,125,125)

            tf.word_wrap = False
            shape.width = Inches(4.5)

# ==============================
# EXPORTAR A PDF
# ==============================
def exportar_pdf(ruta_pptx, ruta_pdf):

    import time

    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = True

    ruta_pptx = os.path.abspath(ruta_pptx)
    ruta_pdf = os.path.abspath(ruta_pdf)

    # esperar a que el archivo exista
    for i in range(10):
        if os.path.exists(ruta_pptx):
            break
        time.sleep(1)

    presentation = powerpoint.Presentations.Open(ruta_pptx, WithWindow=False)

    presentation.SaveAs(ruta_pdf, 32)

    presentation.Close()
    powerpoint.Quit()

# ==============================
# GENERAR INFORME
# ==============================
def generar_informe():

    turno, hora_inicio, hora_fin = obtener_turno()

    global CARPETA_IMAGENES

    # Seleccionar plantilla y carpeta de imágenes según turno
    if turno == "08_AM":
        ruta_plantilla = RUTA_PLANTILLA_8AM
        CARPETA_IMAGENES = CARPETA_IMG_8AM

    elif turno == "02_PM":
        ruta_plantilla = RUTA_PLANTILLA_2PM
        CARPETA_IMAGENES = CARPETA_IMG_2PM

    else:
        ruta_plantilla = RUTA_PLANTILLA_8PM
        CARPETA_IMAGENES = CARPETA_IMG_8PM


    prs = Presentation(ruta_plantilla)

    actualizar_texto(prs, hora_inicio, hora_fin)

    disponibles = obtener_disponibles_portapapeles()

    for agencia in AGENCIAS_ESPERADAS:

        if agencia not in disponibles:

            print("⚠ Falta jefe de sala:", agencia)

    insertar_nombres(prs, disponibles)

    total = procesar_imagenes(prs)

    ahora = datetime.now()

    dia = ahora.day
    mes = MESES[ahora.month]

    nombre_archivo = f"Reporte Seguimiento Operación NUSE 123 {dia} de {mes} {turno}"

    carpeta_informe = os.path.join(CARPETA_SALIDA,"informe")
    carpeta_pdf = os.path.join(CARPETA_SALIDA,f"{mes}-pdf")

    os.makedirs(carpeta_informe,exist_ok=True)
    os.makedirs(carpeta_pdf,exist_ok=True)

    ruta_pptx = os.path.join(carpeta_informe,f"{nombre_archivo}.pptx")
    ruta_pdf = os.path.join(carpeta_pdf,f"{nombre_archivo}.pdf")

    # GUARDAR EL POWERPOINT
    prs.save(ruta_pptx)

    print("Guardando PPT:", ruta_pptx)

    time.sleep(3)

    # EXPORTAR A PDF
    exportar_pdf(ruta_pptx,ruta_pdf)

    print("\nInforme generado correctamente")
    print("Imágenes insertadas:",total)
    print("PPT:",ruta_pptx)
    print("PDF:",ruta_pdf)

    os.startfile(ruta_pdf)

    # abrir el PDF automáticamente
   
    os.startfile(ruta_pdf)

# ==============================
# EJECUTAR
# ==============================
generar_informe()