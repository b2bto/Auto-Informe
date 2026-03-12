import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from PIL import Image
import re

# ==============================
# CONFIGURACION
# ==============================
RUTA_PLANTILLA = r"C:/Users/LuisAlvaroRojasRinco/Documents/Luis-R/Autom/AutInforme-py/plantilla-8.pptx"
CARPETA_IMAGENES = r"C:/Users/LuisAlvaroRojasRinco/Documents/Luis-R/Autom/AutInforme-py/ImagenesInforme/IMG-PNG"
CARPETA_SALIDA = r"C:/Users/LuisAlvaroRojasRinco/Documents/Luis-R/Autom/AutInforme-py/salida"
ARCHIVO_DISPONIBLES = r"C:/Users/LuisAlvaroRojasRinco/Documents/Luis-R/Autom/AutInforme-py/disponibles.txt"

# ==============================
# ACTUALIZAR TEXTO DIAPOSITIVA 1
# ==============================
def actualizar_texto(prs):
    slide = prs.slides[0]
    meses = {1:"enero",2:"febrero",3:"marzo",4:"abril",5:"mayo",6:"junio",7:"julio",8:"agosto",9:"septiembre",10:"octubre",11:"noviembre",12:"diciembre"}
    ahora = datetime.now()
    fecha = f"{ahora.day} de {meses[ahora.month]} del {ahora.year}"
    texto = f"Bogotá D.C. {fecha}\nDesde las 00:00 hasta las 08:00"
    for shape in slide.shapes:
        if shape.name == "txtFecha":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = texto
            p.font.name = "Calibri"
            p.font.size = Pt(24)
            p.font.bold = False
            p.font.color.rgb = RGBColor(255,255,255)

# ==============================
# FORMATEAR NOMBRE
# ==============================
def formatear_nombre(nombre):
    nombre = nombre.strip().lower()
    partes = nombre.split()
    partes = [p.capitalize() for p in partes]
    return " ".join(partes)

# ==============================
# EXTRAER NOMBRES DESDE ARCHIVO
# ==============================
def extraer_nombres_archivo(ruta_archivo):
    nombres_agencias = {
        "SUR": "", "CRUE": "", "MOVILIDAD": "", "IDIGER": "", "BOMBEROS": "", "MEBOG": ""
    }
    if not os.path.exists(ruta_archivo):
        print(f"No se encontró el archivo {ruta_archivo}")
        return nombres_agencias

    with open(ruta_archivo, "r", encoding="utf-8") as f:
        texto = f.read()

    # Buscar todos los pares UBICACIÓN / JEFE DE SALA
    matches = re.findall(r"UBICACIÓN:\s*([A-Z]+).*?JEFE DE SALA:\s*([^\n\r]+)", texto, flags=re.IGNORECASE | re.DOTALL)
    for ubicacion, jefe in matches:
        ubicacion = ubicacion.strip().upper()
        jefe = formatear_nombre(jefe)
        if ubicacion in nombres_agencias:
            nombres_agencias[ubicacion] = jefe

    return nombres_agencias

# ==============================
# INSERTAR NOMBRES EN DIAPOSITIVA
# ==============================
def insertar_nombres(prs, nombres_agencias):
    slide = prs.slides[5]  # Diapositiva donde van los nombres
    for shape in slide.shapes:
        nombre_shape = shape.name.upper()
        if nombre_shape in nombres_agencias and nombres_agencias[nombre_shape]:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = nombres_agencias[nombre_shape]
            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(126,126,126)

# ==============================
# AJUSTAR IMAGEN SIN DEFORMAR
# ==============================
def ajustar_imagen(ruta, marco):
    img = Image.open(ruta)
    img_w, img_h = img.size
    marco_w = marco.width
    marco_h = marco.height
    ratio_img = img_w / img_h
    ratio_marco = marco_w / marco_h
    if ratio_img > ratio_marco:
        new_w = marco_w
        new_h = marco_w / ratio_img
    else:
        new_h = marco_h
        new_w = marco_h * ratio_img
    left = marco.left + (marco_w - new_w) / 2
    top = marco.top + (marco_h - new_h) / 2
    return left, top, new_w, new_h

# ==============================
# PROCESAR IMAGENES
# ==============================
def procesar_imagenes(prs):
    total = 0
    extensiones_validas = (".png", ".jpg", ".jpeg", ".JPG", ".PNG")
    for archivo in os.listdir(CARPETA_IMAGENES):
        if not archivo.endswith(extensiones_validas):
            continue
        if "-D" not in archivo:
            continue
        ruta = os.path.join(CARPETA_IMAGENES, archivo)
        nombre = archivo.lower().replace(".png","").replace(".jpg","").replace(".jpeg","")
        partes = nombre.split("-")
        if len(partes) == 2:
            codigo = partes[1].upper()
        elif len(partes) >= 3:
            codigo = f"{partes[1].upper()}-{partes[2]}"
        else:
            continue
        nombre_marco = f"imgMarco-{codigo}"
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.name.startswith(nombre_marco):
                    marco = shape
                    for s in slide.shapes:
                        if s.name == f"imgAuto-{codigo}":
                            slide.shapes._spTree.remove(s._element)
                    left, top, new_w, new_h = ajustar_imagen(ruta, marco)
                    pic = slide.shapes.add_picture(ruta, left, top, width=new_w, height=new_h)
                    pic.name = f"imgAuto-{codigo}"
                    total += 1
    return total

# ==============================
# GENERAR INFORME
# ==============================
def generar_informe():
    prs = Presentation(RUTA_PLANTILLA)
    actualizar_texto(prs)
    nombres_agencias = extraer_nombres_archivo(ARCHIVO_DISPONIBLES)
    insertar_nombres(prs, nombres_agencias)
    total = procesar_imagenes(prs)
    if not os.path.exists(CARPETA_SALIDA):
        os.makedirs(CARPETA_SALIDA)
    nombre = f"Informe_{datetime.now().strftime('%Y-%m-%d_%H%M%S')}.pptx"
    ruta_final = os.path.join(CARPETA_SALIDA, nombre)
    prs.save(ruta_final)
    print("\nInforme generado correctamente.")
    print("Total imágenes insertadas:", total)
    print("Archivo guardado en:", ruta_final)

# ==============================
# EJECUTAR
# ==============================
generar_informe()