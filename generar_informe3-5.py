import os
import time
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import pyperclip
import win32com.client

# ==============================
# CONFIGURACION
# ==============================
RUTA_PLANTILLA = r"C:/Users/LuisAlvaroRojasRinco/Documents/Luis-R/Autom/AutInforme-py/plantilla-8am.pptx"
CARPETA_IMAGENES = r"C:/Users/LuisAlvaroRojasRinco/Documents/Luis-R/Autom/AutInforme-py/ImagenesInforme/IMG-PNG"
CARPETA_SALIDA = r"C:/Users/LuisAlvaroRojasRinco/Documents/Luis-R/Autom/AutInforme-py/salida"

AGENCIAS_ESPERADAS = [
    "SUR",
    "CRUE",
    "MOVILIDAD",
    "BOMBEROS",
    "MEBOG",
    "IDIGER"
]

# ==============================
# MESES
# ==============================
MESES = {
    1:"enero",2:"febrero",3:"marzo",4:"abril",
    5:"mayo",6:"junio",7:"julio",8:"agosto",
    9:"septiembre",10:"octubre",11:"noviembre",12:"diciembre"
}

# ==============================
# ACTUALIZAR TEXTO DIAPOSITIVA 1
# ==============================
def actualizar_texto(prs):

    slide = prs.slides[0]

    ahora = datetime.now()
    fecha = f"{ahora.day} de {MESES[ahora.month]} del {ahora.year}"

    texto = f"Bogotá D.C. {fecha}\nDesde las 00:00 hasta las 08:00"

    for shape in slide.shapes:

        if shape.name == "txtFecha":

            tf = shape.text_frame
            tf.clear()

            p = tf.paragraphs[0]
            p.text = texto

            p.font.name = "Calibri"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255,255,255)

            p.alignment = PP_ALIGN.CENTER

# ==============================
# FORMATEAR NOMBRES
# ==============================
def formatear_nombre(nombre):

    partes = nombre.strip().lower().split()
    partes = [p.capitalize() for p in partes]

    return " ".join(partes)

# ==============================
# EXTRAER NOMBRES WHATSAPP
# ==============================
def obtener_disponibles_portapapeles():

    disponibles = {}

    texto = pyperclip.paste()
    lineas = texto.splitlines()

    agencia = ""
    nombre = ""

    for linea in lineas:

        linea = linea.strip()

        if linea.upper().startswith("UBICACIÓN:"):

            agencia = linea.split(":")[1].strip().upper()

        elif linea.upper().startswith("JEFE DE SALA:"):

            nombre = formatear_nombre(linea.split(":")[1].strip())

            if agencia and nombre:

                disponibles[agencia] = nombre
                agencia = ""
                nombre = ""

        elif "-" in linea:

            partes = linea.split("-")

            if len(partes) == 2:

                agencia = partes[0].strip().upper()
                nombre = formatear_nombre(partes[1].strip())

                disponibles[agencia] = nombre

    print("Disponibles detectados:", disponibles)

    return disponibles

# ==============================
# AJUSTAR IMAGEN
# ==============================
def ajustar_imagen(ruta, marco):

    img = Image.open(ruta)
    img_w, img_h = img.size

    marco_w = marco.width
    marco_h = marco.height

    ratio_img = img_w / img_h
    ratio_marco = marco_w / marco_h

    if ratio_img > ratio_marco:

        new_w = marco_w
        new_h = marco_w / ratio_img

    else:

        new_h = marco_h
        new_w = marco_h * ratio_img

    left = marco.left + (marco_w - new_w) / 2
    top = marco.top + (marco_h - new_h) / 2

    return left, top, new_w, new_h

# ==============================
# PROCESAR IMAGENES (MEJORADO)
# ==============================
def procesar_imagenes(prs):

    total = 0
    extensiones = (".png",".jpg",".jpeg")

    for archivo in os.listdir(CARPETA_IMAGENES):

        if not archivo.lower().endswith(extensiones):
            continue

        nombre = archivo.lower()

        if "-d" not in nombre:
            continue

        ruta = os.path.join(CARPETA_IMAGENES, archivo)

        codigo = nombre.split("-d")[1].split(".")[0].upper()
        codigo = f"D{codigo}"

        nombre_marco = f"imgMarco-{codigo}"

        marco_encontrado = False

        for slide in prs.slides:

            if marco_encontrado:
                break

            for shape in slide.shapes:

                if shape.name.upper().startswith(nombre_marco.upper()):

                    marco_encontrado = True

                    for s in slide.shapes:

                        if s.name == f"imgAuto-{codigo}":

                            slide.shapes._spTree.remove(s._element)

                    left, top, new_w, new_h = ajustar_imagen(ruta, shape)

                    pic = slide.shapes.add_picture(
                        ruta,
                        left,
                        top,
                        width=new_w,
                        height=new_h
                    )

                    pic.name = f"imgAuto-{codigo}"

                    total += 1

                    print("Imagen insertada:", archivo)

                    break

        if not marco_encontrado:

            print("⚠ No se encontró marco para:", archivo)

    return total

# ==============================
# LIMPIAR RANGO
# ==============================
def limpiar_rango(nombre):

    rangos = ["mayor","teniente","sargento","capitan","coronel","subteniente"]

    partes = nombre.lower().split()

    partes = [p for p in partes if p not in rangos]

    return " ".join(partes).title()

# ==============================
# INSERTAR NOMBRES
# ==============================
def insertar_nombres(prs, disponibles):

    slide = prs.slides[2]

    for shape in slide.shapes:

        agencia = shape.name.strip().upper()

        if agencia in AGENCIAS_ESPERADAS and shape.has_text_frame:

            if agencia in disponibles:
                nombre = limpiar_rango(disponibles[agencia])
            else:
                nombre = "Sin reporte"

            tf = shape.text_frame
            tf.clear()

            p = tf.paragraphs[0]

            p.text = f"{agencia}: {nombre}"

            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(125,125,125)

            tf.word_wrap = False
            shape.width = Inches(4.5)

# ==============================
# EXPORTAR A PDF
# ==============================
def exportar_pdf(ruta_pptx, ruta_pdf):

    import time

    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = True

    ruta_pptx = os.path.abspath(ruta_pptx)
    ruta_pdf = os.path.abspath(ruta_pdf)

    # esperar a que el archivo exista
    for i in range(10):
        if os.path.exists(ruta_pptx):
            break
        time.sleep(1)

    presentation = powerpoint.Presentations.Open(ruta_pptx, WithWindow=False)

    presentation.SaveAs(ruta_pdf, 32)

    presentation.Close()
    powerpoint.Quit()

# ==============================
# GENERAR INFORME
# ==============================
def generar_informe():

    prs = Presentation(RUTA_PLANTILLA)

    actualizar_texto(prs)

    disponibles = obtener_disponibles_portapapeles()

    for agencia in AGENCIAS_ESPERADAS:

        if agencia not in disponibles:

            print("⚠ Falta jefe de sala:", agencia)

    insertar_nombres(prs, disponibles)

    total = procesar_imagenes(prs)

    ahora = datetime.now()

    dia = ahora.day
    mes = MESES[ahora.month]

    nombre_archivo = f"Reporte Seguimiento Operación NUSE 123 {dia} de {mes} 08_AM"

    carpeta_informe = os.path.join(CARPETA_SALIDA,"informe")
    carpeta_pdf = os.path.join(CARPETA_SALIDA,f"{mes}-pdf")

    os.makedirs(carpeta_informe,exist_ok=True)
    os.makedirs(carpeta_pdf,exist_ok=True)

    ruta_pptx = os.path.join(carpeta_informe,f"{nombre_archivo}.pptx")
    ruta_pdf = os.path.join(carpeta_pdf,f"{nombre_archivo}.pdf")

    # GUARDAR EL POWERPOINT
    prs.save(ruta_pptx)

    print("Guardando PPT:", ruta_pptx)

    # esperar a que el archivo termine de guardarse
    import time
    time.sleep(3)

    # EXPORTAR A PDF
    exportar_pdf(ruta_pptx,ruta_pdf)

    print("\nInforme generado correctamente")
    print("Imágenes insertadas:",total)
    print("PPT:",ruta_pptx)
    print("PDF:",ruta_pdf)

    # abrir el PDF automáticamente
   
    os.startfile(ruta_pdf)

# ==============================
# EJECUTAR
# ==============================
generar_informe()
