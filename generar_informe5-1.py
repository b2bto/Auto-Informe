# Librerías estándar
import os
import sys
import copy
import time
from datetime import datetime, timedelta

# Librerías externas
import pyperclip
import win32com.client
from PIL import Image, ImageEnhance
from openpyxl import load_workbook, Workbook

# PowerPoint
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE


# ==============================
# RUTA BASE DEL PROGRAMA
# ==============================
if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))
RUTA_EXCEL = os.path.join(BASE_DIR, "incidentes_mes.xlsx")
RUTA_PLANTILLA_8AM = os.path.join(BASE_DIR,"plantilla-8am.pptx")
RUTA_PLANTILLA_2PM = os.path.join(BASE_DIR,"plantilla-2pm.pptx")
RUTA_PLANTILLA_8PM = os.path.join(BASE_DIR,"plantilla-8pm.pptx")

#Rutas de imágenes automáticas
CARPETA_IMG_8AM = os.path.join(BASE_DIR,"ImagenesInforme","IMG-8AM")
CARPETA_IMG_2PM = os.path.join(BASE_DIR,"ImagenesInforme","IMG-2PM")
CARPETA_IMG_8PM = os.path.join(BASE_DIR,"ImagenesInforme","IMG-8PM")

#Carpeta de salida
CARPETA_SALIDA = os.path.join(BASE_DIR,"salida")

AGENCIAS_ESPERADAS = [
    "SUR",
    "CRUE",
    "MOVILIDAD",
    "BOMBEROS",
    "MEBOG",
    "IDIGER"
]

# ==============================
# MESES
# ==============================
MESES = {
    1:"enero",2:"febrero",3:"marzo",4:"abril",
    5:"mayo",6:"junio",7:"julio",8:"agosto",
    9:"septiembre",10:"octubre",11:"noviembre",12:"diciembre"
}

# ==============================
# FORMATO CORPORATIVO TEXTO
# ==============================
def aplicar_formato(run, size=16, r=126, g=126, b=126):
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor(r,g,b)

# ==============================
# CLONAR DIAPOSITIVA (CON POSICIÓN)
# ==============================
def clonar_diapositiva(prs, slide_index, posicion_destino):

    slide = prs.slides[slide_index]
    layout = slide.slide_layout

    new_slide = prs.slides.add_slide(layout)

    for shape in slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # MOVER LA SLIDE A LA POSICIÓN CORRECTA
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)

    slide_id = slides[-1]  # la última (recién creada)

    slide_id_list.remove(slide_id)
    slide_id_list.insert(posicion_destino, slide_id)

    return new_slide

# DETECTAR TURNO AUTOMATICO
# ==============================
def obtener_turno():

    ahora = datetime.now()
    hora = ahora.hour
    minuto = ahora.minute

    hora_decimal = hora + minuto/60

    if hora_decimal < 8.5:
        return "08_AM", "00:00", "08:00"

    elif hora_decimal < 14.5:
        return "02_PM", "00:00", "14:00"

    elif hora_decimal < 20.5:
        return "08_PM", "00:00", "20:00"

    else:
        return "08_PM", "00:00", "20:00"

# ACTUALIZAR TEXTO DIAPOSITIVA 1
# ==============================
def actualizar_texto(prs, hora_inicio, hora_fin):

    slide = prs.slides[0]

    ahora = datetime.now()
    dia = f"{ahora.day:02d}"
    mes = MESES[ahora.month].capitalize()
    anio = ahora.year

    linea1 = f"Bogotá D.C. {dia} de {mes} del {anio}"
    linea2 = f"Desde las {hora_inicio} hasta las {hora_fin}"

    for shape in slide.shapes:

        if shape.name == "txtFecha" and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()
            
            p1 = tf.paragraphs[0]
            p1.text = linea1
            p1.font.name = "Calibri"
            p1.font.size = Pt(24)
            p1.font.color.rgb = RGBColor(255,255,255)
            p1.alignment = PP_ALIGN.CENTER

            p2 = tf.add_paragraph()
            p2.text = linea2
            p2.font.name = "Calibri"
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(255,255,255)
            p2.alignment = PP_ALIGN.CENTER

# ==============================
# ACTUALIZAR TEXTO DIAPOSITIVA 16 (solo 8AM)
# ==============================

def actualizar_estadisticas_8am(prs):

    slide = prs.slides[15]  # diapositiva 16

    # CALCULAR FECHAS
    hoy = datetime.now()
    ayer = hoy - timedelta(days=1)
    antier = hoy - timedelta(days=2)

    fecha_ayer = ayer.strftime("%d/%m/%Y")
    fecha_antier = antier.strftime("%d/%m/%Y")

    # TEXTOS
    txt1 = f"Estadística general de llamadas con corte desde las 00:00 AM hasta las 11:59 PM del {fecha_ayer}"
    txt2 = f"Estadística total llamadas con corte desde las 00:00 AM hasta las 11:59 PM del {fecha_antier} y {fecha_ayer}"

    for shape in slide.shapes:

        if shape.name == "txt-D16-1" and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()

            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            run = p.add_run()
            run.text = txt1

            aplicar_formato(run)

        elif shape.name == "txt-D16-2" and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()

            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            run = p.add_run()
            run.text = txt2

            aplicar_formato(run)

# ==============================
# LEER INCIDENTES DESDE EXCEL
# ==============================
def leer_incidentes():
    if not os.path.exists(RUTA_EXCEL):
        return []

    try:
        wb = load_workbook(RUTA_EXCEL, data_only=True)
    except PermissionError:
        print(f"❌ ERROR: El archivo Excel está abierto: {RUTA_EXCEL}")
        print("Por favor, ciérralo e intenta de nuevo.")
        time.sleep(5) 
        sys.exit(1) 
    except Exception as e:
        print(f"❌ Error inesperado al abrir el Excel: {e}")
        return []

    if "Incidentes" not in wb.sheetnames:
        return []

    ws = wb["Incidentes"]
    incidentes = []

    for row in ws.iter_rows(min_row=2, values_only=True):
        if not row or len(row) < 3:
            continue
            
        codigo = row[0]
        agencia = row[1]
        tipo = row[2]

        if codigo and agencia and tipo:
            incidentes.append({
                "codigo": str(codigo).strip(),
                "agencia": str(agencia).strip().upper(),
                "tipo": str(tipo).strip().lower()
            })
    return incidentes

# ==============================
# CREAR EXCEL DE INCIDENTES SI NO EXISTE
# ==============================
def crear_excel_incidentes():
    if os.path.exists(RUTA_EXCEL):
        return
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Incidentes"
    ws1.append(["Codigo", "Agencia", "Tipo"])

    ws2 = wb.create_sheet("Novedades")
    ws2.append(["Hora", "Agencia", "Novedad", "Estado"]) 
    wb.save(RUTA_EXCEL)
    print("Excel creado:", RUTA_EXCEL)

# ==============================
# LEER NOVEDADES DESDE EXCEL
# ==============================
def leer_novedades():
    ruta = os.path.join(BASE_DIR, RUTA_EXCEL)

    if not os.path.exists(ruta):
        return []

    wb = load_workbook(ruta, data_only=True)

    if "Novedades" not in wb.sheetnames:
        return []

    ws = wb["Novedades"]
    novedades = []

    for row in ws.iter_rows(min_row=2, values_only=True):
        
        # --- AQUÍ VA LA VALIDACIÓN DE SEGURIDAD ---
        if not row or len(row) < 4:
            continue  # Si la fila está vacía o le faltan columnas, salta a la siguiente
        # ------------------------------------------

        fecha = row[0]
        agencia = row[1]
        novedad = row[2]
        estado = row[3]

        if novedad is None:
            continue

        if estado is None:
            continue

        estado = str(estado).strip().upper()

        if estado != "ACTIVA":
            continue

        texto = str(novedad).strip()

        if texto:
            novedades.append(texto)

    return novedades

# ==============================
# CONTAR TIPOS DE INCIDENTES
# ==============================
def contar_tipos_incidente(incidentes):
    malware = 0
    ransomware = 0
    otros = 0

    for inc in incidentes:
        tipo = inc["tipo"].lower().strip() # Normalización total
        if "malware" in tipo:
            malware += 1
        elif "ransomware" in tipo:
            ransomware += 1
        else:
            otros += 1
    return malware, ransomware, otros


# ==============================
# DETECTAR AMENAZA MAS REPETIDA
# ==============================
def amenaza_mas_recurrente(incidentes):

    conteo = {}

    for inc in incidentes:

        codigo = inc["codigo"]

        if codigo not in conteo:
            conteo[codigo] = 0

        conteo[codigo] += 1

    if not conteo:
        return None, 0

    codigo_top = max(conteo, key=conteo.get)
    cantidad = conteo[codigo_top]

    return codigo_top, cantidad

# ==============================
# DETECTAR AGENCIAS AFECTADAS
# ==============================
def detectar_agencias(incidentes):

    mapa = {
        "SUR": "S.U.R.",
        "CRUE": "CRUE",
        "CAD": "CAD"
    }

    agencias = set()

    for inc in incidentes:

        codigo = inc["agencia"]

        if codigo in mapa:
            agencias.add(mapa[codigo])

    return sorted(list(agencias))

# ==============================
# FORMATEAR LISTA AGENCIAS
# ==============================
def formatear_agencias(lista):

    if not lista:
        return ""

    if len(lista) == 1:
        return lista[0]

    if len(lista) == 2:
        return f"{lista[0]} y {lista[1]}"

    return ", ".join(lista[:-1]) + f" y {lista[-1]}"

# ==============================
# GENERAR TEXTO INCIDENTES D7
# ==============================
def generar_texto_incidentes():
    
    incidentes = leer_incidentes()

    malware, ransomware, otros = contar_tipos_incidente(incidentes)
    codigo_top, cantidad_top = amenaza_mas_recurrente(incidentes)
    tipos = []

    if malware > 0:
        tipos.append(f"{malware} eventos tipo malware")

    if ransomware > 0:
        tipos.append(f"{ransomware} eventos tipo ransomware")

    if otros > 0:
        tipos.append(f"{otros} eventos de otro tipo")

    tipos_txt = " y ".join(tipos)

    frase_top = ""

    if cantidad_top > 1:

        frase_top = (
            f" Adicionalmente, se identifica que el evento más recurrente fue "
            f"{codigo_top} con {cantidad_top} registros durante el período analizado."
        )
    

    hoy = datetime.now()
    mes = MESES[hoy.month]
    anio = hoy.year

    if not incidentes:

        return (
        f"El registro del top de amenazas corresponde al período comprendido "
        f"entre el 1 de {mes} de {anio} y la fecha registrada en el pantallazo. "
        f"En el período indicado no se han identificado incidentes de seguridad "
        f"asociados a software malicioso en las plataformas monitoreadas."
        )

    total = len(incidentes)

    agencias = detectar_agencias(incidentes)
    agencias_txt = formatear_agencias(agencias)

    codigos = []

    for inc in incidentes:
        codigo = inc["codigo"]
        tipo = inc["tipo"].capitalize()
        codigos.append(f"({codigo} {tipo})")

    if len(codigos) == 1:
        lista = codigos[0]
    else:
        lista = ", ".join(codigos[:-1]) + " y " + codigos[-1]

    texto = (
        f"El registro del top de amenazas corresponde al período comprendido "
        f"entre el 1 de {mes} de {anio} y la fecha registrada en el pantallazo. "
        f"En el período indicado se han identificado {total} incidentes de seguridad "
        f"en las agencias {agencias_txt}, correspondientes a {tipos_txt}, "
        f"específicamente: {lista}. "
        f"Es importante reforzar los controles de seguridad asociados al software malicioso, "
        f"tales como: monitoreo periódico de eventos de seguridad, actualización continua "
        f"de herramientas antivirus, respaldo de las soluciones de seguridad y programas "
        f"de concientización a los usuarios sobre phishing, descargas no autorizadas y "
        f"buenas prácticas de seguridad de la información. La consolidación de estas "
        f"medidas permitirá reducir el riesgo de incidentes, mantener la resiliencia "
        f"operativa y garantizar el cumplimiento de los lineamientos de la operación."
        f"{frase_top}"
    )

    return texto

# ==============================
# ACTUALIZAR DIAPOSITIVA 7
# ==============================
def actualizar_diapositiva_7(prs):

    slide = prs.slides[6]

    texto = generar_texto_incidentes()

    for shape in slide.shapes:

        if shape.name == "txt-D7" and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()

            # CONFIGURACION DEL CUADRO DE TEXTO
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.JUSTIFY   # TEXTO JUSTIFICADO

            run = p.add_run()
            run.text = texto

            aplicar_formato(run, size=14)   # TAMAÑO DE LETRA 14

# ==============================
# ACTUALIZAR TITULO NOVEDADES
# ==============================
def actualizar_titulo_novedades(slide, numero):

    for shape in slide.shapes:

        if shape.name.startswith("titulo-D") and shape.has_text_frame:

            tf = shape.text_frame
            tf.clear()

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT

            run = p.add_run()
            run.text = f"{numero}. Novedades"

            run.font.name = "Calibri"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(27, 95, 167)

# ==============================
# CALCULAR TAMAÑO DINÁMICO TEXTO
# ==============================
def calcular_tamano_texto(texto):

    longitud = len(texto)

    if longitud < 80:
        return 18
    elif longitud < 150:
        return 16
    elif longitud < 250:
        return 14
    else:
        return 12
    
# ==============================
# ACTUALIZAR DIAPOSITIVA 26
# ==============================
def actualizar_diapositiva_26(prs):

    slide_index = 25  # diapositiva base (26)

    # Validación
    if len(prs.slides) <= slide_index:
        print("❌ La plantilla no tiene suficientes diapositivas")
        return

    novedades = leer_novedades()

    if not novedades:
        novedades = [
            "Durante el periodo del informe no se registran novedades relevantes en la operación."
        ]

    MAX_POR_SLIDE = 3
    BASE_TITULO = 15

    # dividir en bloques
    bloques = [
        novedades[i:i + MAX_POR_SLIDE]
        for i in range(0, len(novedades), MAX_POR_SLIDE)
    ]

    for i, bloque in enumerate(bloques):

        # ==============================
        # CREAR / OBTENER SLIDE
        # ==============================
        if i == 0:
            slide = prs.slides[slide_index]
        else:
            posicion = len(prs.slides) - 2  # 🔥 seguro
            slide = clonar_diapositiva(prs, slide_index, posicion)

        # ==============================
        # RENOMBRAR SHAPES (CLAVE 🔥)
        # ==============================
        for shp in slide.shapes:

            if shp.name == "txt-D26":
                shp.name = f"txt-D{26 + i}"

            if shp.name == "titulo-D26":
                shp.name = f"titulo-D{26 + i}"

        # ==============================
        # ACTUALIZAR TÍTULO
        # ==============================
        for shape in slide.shapes:

            if shape.name.startswith("titulo-D") and shape.has_text_frame:

                tf = shape.text_frame
                tf.clear()

                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT

                run = p.add_run()
                run.text = f"{BASE_TITULO + i}. Novedades"

                run.font.name = "Calibri"
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = RGBColor(27, 95, 167)

        # ==============================
        # ESCRIBIR NOVEDADES
        # ==============================
        nombre_shape = f"txt-D{26 + i}"

        for shape in slide.shapes:

            if shape.name == nombre_shape and shape.has_text_frame:

                tf = shape.text_frame
                tf.clear()

                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.NONE  # 🔥 no romper formato

                for j, novedad in enumerate(bloque):

                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(10)

                    # 🔥 tamaño dinámico
                    tamano = calcular_tamano_texto(novedad)

                    run = p.add_run()
                    run.text = f"• {novedad}"

                    run.font.name = "Calibri"
                    run.font.size = Pt(tamano)
                    run.font.color.rgb = RGBColor(156, 156, 156)
# ==============================
# FORMATEAR NOMBRES
# ==============================
def formatear_nombre(nombre):

    partes = nombre.strip().lower().split()
    partes = [p.capitalize() for p in partes]

    return " ".join(partes)

# ==============================
# EXTRAER NOMBRES WHATSAPP
# ==============================
def obtener_disponibles_portapapeles():

    disponibles = {}

    texto = pyperclip.paste()
    lineas = texto.splitlines()

    agencia = ""
    nombre = ""

    for linea in lineas:

        linea = linea.strip()

        if linea.upper().startswith("UBICACIÓN:"):

            agencia = linea.split(":")[1].strip().upper()

        elif linea.upper().startswith("JEFE DE SALA:"):

            nombre = formatear_nombre(linea.split(":")[1].strip())

            if agencia and nombre:

                disponibles[agencia] = nombre
                agencia = ""
                nombre = ""

        elif "-" in linea:

            partes = linea.split("-")

            if len(partes) == 2:

                agencia = partes[0].strip().upper()
                nombre = formatear_nombre(partes[1].strip())

                disponibles[agencia] = nombre

    print("Disponibles detectados:", disponibles)

    return disponibles

# ==========================================================
# AJUSTAR IMAGEN CON MARGEN Y MEJORAR CALIDAD (PRO)
# ==========================================================
def ajustar_imagen_con_margen(ruta, marco, color_fondo=(255, 255, 255)):

    try:
        with Image.open(ruta) as img_raw:
            img = img_raw.convert("RGB")
            img_w, img_h = img.size

            marco_w = int(marco.width)
            marco_h = int(marco.height)

            MAX_PIX = 4000
            escala_w = min(marco_w, MAX_PIX)
            escala_h = min(marco_h, MAX_PIX)

            ratio_img = img_w / img_h
            ratio_marco = escala_w / escala_h

            if ratio_img > ratio_marco:
                new_w = escala_w
                new_h = int(escala_w / ratio_img)
            else:
                new_h = escala_h
                new_w = int(escala_h * ratio_img)

            img_resized = img.resize((new_w, new_h), Image.LANCZOS)

            sharp = ImageEnhance.Sharpness(img_resized)
            img_resized = sharp.enhance(1.2)

            fondo = Image.new("RGB", (new_w, new_h), color_fondo)

            pos_x = (new_w - img_resized.width) // 2
            pos_y = (new_h - img_resized.height) // 2
            fondo.paste(img_resized, (pos_x, pos_y))

            nombre_base = os.path.splitext(os.path.basename(ruta))[0]
            ruta_temp = os.path.join(os.path.dirname(ruta), f"{nombre_base}_PRO.png")

            # 🔥 guardar SIEMPRE
            fondo.save(ruta_temp, format="PNG", dpi=(300, 300))

            return ruta_temp, marco.left, marco.top, marco.width, marco.height

    except Exception as e:
        print(f"❌ Error procesando imagen {ruta}: {e}")

        # 🔥 fallback: usar imagen original (para que NO se rompa el flujo)
        return ruta, marco.left, marco.top, marco.width, marco.height

# ==============================
# PROCESAR IMAGENES (MEJORADO)
# ==============================
def procesar_imagenes(prs, CARPETA_IMAGENES):

    # 🔥 VALIDACIÓN SEGURA
    if not os.path.exists(CARPETA_IMAGENES):
        print("⚠ Carpeta de imágenes no encontrada:", CARPETA_IMAGENES)
        return 0

    total = 0
    extensiones = (".png",".jpg",".jpeg")

    #ruta_temp = os.path.join(CARPETA_IMAGENES, "_temp_img.png")

    for archivo in os.listdir(CARPETA_IMAGENES):

        if not archivo.lower().endswith(extensiones):
            continue

        nombre = archivo.lower()

        if "-d" not in nombre:
            continue

        ruta = os.path.join(CARPETA_IMAGENES, archivo)

        codigo = nombre.split("-d")[1].split(".")[0].upper()
        codigo = f"D{codigo}"

        nombre_marco = f"imgMarco-{codigo}"

        marco_encontrado = False

        for slide in prs.slides:

            if marco_encontrado:
                break

            for shape in slide.shapes:

                if shape.name.upper().startswith(nombre_marco.upper()):

                    marco_encontrado = True

                    # eliminar imagen anterior si existe
                    for s in slide.shapes:
                        if s.name == f"imgAuto-{codigo}":
                            slide.shapes._spTree.remove(s._element)

                    # generar imagen ajustada con márgenes
                    # ruta_temp, left, top, new_w, new_h = ajustar_imagen_con_margen(ruta, shape, color_fondo=(255,255,255))

                    ruta_temp, left, top, new_w, new_h = ajustar_imagen_con_margen(ruta, shape, color_fondo=(255,255,255))

                    # insertar imagen directamente
                    pic = slide.shapes.add_picture(
                        ruta_temp,
                        left,
                        top,
                        width=new_w,
                        height=new_h
                    )
                   
                    pic.name = f"imgAuto-{codigo}"

                    total += 1

                    print("Imagen insertada:", archivo)

                    break

        if not marco_encontrado:
            print("⚠ No se encontró marco para:", archivo)

    # limpiar imagen temporal

    return total
# ==============================
# LIMPIAR IMAGENES TEMPORALES
# ==============================
def limpiar_imagenes_temporales(carpeta):

    if not os.path.exists(carpeta):
        return

    for archivo in os.listdir(carpeta):

        if archivo.endswith("_PRO.png"):

            ruta = os.path.join(carpeta, archivo)

            try:
                os.remove(ruta)
                print("🧹 Eliminado:", archivo)
            except Exception as e:
                print(f"⚠ No se pudo eliminar {archivo}: {e}")

# ==============================
# LIMPIAR RANGO
# ==============================
def limpiar_rango(nombre):

    rangos = [
        "mayor","my","m.y",
        "teniente","tn","t.te",
        "capitan","cap","cpt",
        "coronel","cr","cor",
        "subteniente","st",
        "sargento","sg","sgto",
        "intendente","int","it",
        "patrullero","pt",
        "oficial","of"
    ]

    partes = nombre.lower().replace(".","").split()

    partes = [p for p in partes if p not in rangos]

    return " ".join(partes).title()

# ==============================
# INSERTAR NOMBRES
# ==============================
def insertar_nombres(prs, disponibles):

    slide = prs.slides[2]

    for shape in slide.shapes:

        agencia = shape.name.strip().upper()

        if agencia in AGENCIAS_ESPERADAS and shape.has_text_frame:

            if agencia in disponibles:
                nombre = limpiar_rango(disponibles[agencia])
            else:
                nombre = "Sin reporte"

            tf = shape.text_frame
            tf.clear()

            p = tf.paragraphs[0]

            run = p.add_run()
            run.text = f"{agencia}: {nombre}"
            aplicar_formato(run)

            tf.word_wrap = False
            shape.width = Inches(4.5)

# ==============================
# EXPORTAR A PDF
# ==============================
def exportar_pdf(ruta_pptx, ruta_pdf):
    powerpoint = None
    try:
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        # Abrir de forma invisible
        presentation = powerpoint.Presentations.Open(os.path.abspath(ruta_pptx), WithWindow=False)
        presentation.SaveAs(os.path.abspath(ruta_pdf), 32)
        presentation.Close()
    except Exception as e:
        print("⚠ No se pudo exportar a PDF.")
        print("Posibles causas:")
        print("- PowerPoint no está instalado")
        print("- Error en COM de Windows")
        print("- Archivo PPT abierto")
        print("Detalle:", e)
    finally:
        if powerpoint:
            powerpoint.Quit()

# ==============================
# GENERAR INFORME
# ==============================
def generar_informe(turno_manual=None):

    crear_excel_incidentes()

    # 1. Definir el turno
    if turno_manual == "mañana":
        turno, hora_inicio, hora_fin = "08_AM", "00:00", "08:00"
    elif turno_manual == "tarde":
        turno, hora_inicio, hora_fin = "02_PM", "00:00", "14:00"
    elif turno_manual == "noche":
        turno, hora_inicio, hora_fin = "08_PM", "00:00", "20:00"
    else:
        turno, hora_inicio, hora_fin = obtener_turno()

    # 2. Seleccionar rutas según el turno
    if turno == "08_AM":
        ruta_plantilla = RUTA_PLANTILLA_8AM
        CARPETA_IMAGENES = CARPETA_IMG_8AM
    elif turno == "02_PM":
        ruta_plantilla = RUTA_PLANTILLA_2PM
        CARPETA_IMAGENES = CARPETA_IMG_2PM
    else:
        ruta_plantilla = RUTA_PLANTILLA_8PM
        CARPETA_IMAGENES = CARPETA_IMG_8PM

    # --- AQUÍ VA LA VALIDACIÓN DE SEGURIDAD ---
    if not os.path.exists(ruta_plantilla):
        print(f"❌ ERROR: No se encontró la plantilla en: {ruta_plantilla}")
        return # Esto detiene la función para que no intente abrir algo que no existe
    # ------------------------------------------

    # 3. Abrir plantilla con protección
    try:
        prs = Presentation(ruta_plantilla)
    except Exception as e:
        print(f"❌ Error al abrir la presentación: {e}")
        return

    actualizar_texto(prs, hora_inicio, hora_fin)
    actualizar_diapositiva_7(prs)
    actualizar_diapositiva_26(prs)

    disponibles = obtener_disponibles_portapapeles()

    for agencia in AGENCIAS_ESPERADAS:
        if agencia not in disponibles:
            print("⚠ Falta jefe de sala:", agencia)

    insertar_nombres(prs, disponibles)

    # PROCESAR IMÁGENES
    total = procesar_imagenes(prs, CARPETA_IMAGENES)
    
    # solo para informe 8AM
    if turno == "08_AM":
        actualizar_estadisticas_8am(prs)

    ahora = datetime.now()
    dia = ahora.day
    mes = MESES[ahora.month]

    nombre_archivo = f"Reporte Seguimiento Operación NUSE 123 {dia} de {mes} {turno}"

    carpeta_informe = os.path.join(CARPETA_SALIDA, "informe")
    carpeta_pdf = os.path.join(CARPETA_SALIDA, f"{mes}-pdf")

    os.makedirs(carpeta_informe, exist_ok=True)
    os.makedirs(carpeta_pdf, exist_ok=True)

    ruta_pptx = os.path.join(carpeta_informe, f"{nombre_archivo}.pptx")
    ruta_pdf = os.path.join(carpeta_pdf, f"{nombre_archivo}.pdf")

    # --- MEJORA AL GUARDAR ---
    try:
        prs.save(ruta_pptx)
        print("Guardando PPT:", ruta_pptx)
    except PermissionError:
        print(f"❌ ERROR: No se pudo guardar el PPT. ¿Está abierto el archivo {nombre_archivo}.pptx?")
        return
    # -------------------------

#    time.sleep(3)

   # EXPORTAR PDF
    exportar_pdf(ruta_pptx, ruta_pdf)

    # 🔥 VALIDAR SI EL PDF SE CREÓ
    if not os.path.exists(ruta_pdf):
        print("\n⚠ No se generó el PDF")
        print("Revisa PowerPoint o posibles errores de exportación")
        print("El archivo PPT sí fue generado correctamente:")
        print("PPT:", ruta_pptx)

        os.startfile(ruta_pptx) # 🔥 abre el PPT automáticamente
        return

    print("\nInforme generado correctamente")
    print("Imágenes insertadas:", total)
    print("PPT:", ruta_pptx)
    print("PDF:", ruta_pdf)

    # Abrir automáticamente el PDF
    os.startfile(ruta_pdf)
    limpiar_imagenes_temporales(CARPETA_IMAGENES)

# ==============================
# EJECUTAR SOLO EN CONSOLA
# ==============================
if __name__ == "__main__":

    print("Modo consola")

    turno, h1, h2 = obtener_turno()
    generar_informe()